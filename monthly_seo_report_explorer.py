"""
monthly_seo_report_explorer.py
──────────────────────────────────────────────────────────────────────────
Standalone Explorer/Bing-only monthly report builder — same
detect -> review -> build architecture as monthly_seo_report_google.py,
kept as its own separate script (not combined with Google) since the
screenshot format, detection approach, and flagging criteria are all
different:

  • Folder layout (per day folder) — files sit DIRECTLY in the day
    folder, no subfolder involved:
        <day>/<keyword ...> pg N no M[,M2&M3...].jpg
        <day>/<keyword ...> no results.jpg   (skipped entirely)

  • Bing prints an "N-M of X results" header instead of Google's blue
    title-link layout carrying no such text — detection reads that
    header via OCR to get the page number and how many results are on
    the page, then locates each individual result by its own blue
    title-link color, the same left-alignment/spacing idea as the
    Google pipeline but tuned for Bing's specific colors/layout.

  • Flagging here is COURT / LEGAL / FINANCIAL-CRIME keyword matching
    (see FLAGGED_KEYWORDS) — NOT name-matching like the Google script,
    since Explorer/Bing results are being screened for adverse-media
    mentions, not confirming the person's name appears. The filename's
    own "no 6,8&10" claim is NOT used for flagging — that comes only
    from what's actually detected in the image.

  • A results page commonly holds more rows than fit legibly on one
    slide, so a screenshot is normally split into 2 output images
    ("Part 1 of 2" / "Part 2 of 2") — but this is driven by the ACTUAL
    detected row count, not a fixed assumption.

  • Same review-before-build workflow as Google: scan_and_detect() finds
    every screenshot and detects rows only (no images built yet), a
    review tool shows every row so a human can correct anything wrong —
    including a real VISUAL crop editor (drag directly on the raw
    screenshot to redraw a row's boundary, add/remove/reorder rows) via
    the companion review_report_explorer.py — and only then does
    build_report_from_manifest() actually produce the final images/CSV/
    PPTX, using exactly what was confirmed.

  • The CSV and PPTX metadata reflect what's actually DETECTED (the
    OCR'd page number, the real flagged rows) — not the filename's own
    "PG#N NO#M" claim, which has no way to catch a stale or simply wrong
    number typed when the file was originally named.

Run it with:
    python monthly_seo_report_explorer.py
"""

import json
import csv
import copy
import os
import queue
import re
import sys
import threading
import traceback
from collections import defaultdict
from concurrent.futures import ThreadPoolExecutor, as_completed
from pathlib import Path

import numpy as np
import pytesseract
from PIL import Image, ImageDraw, ImageEnhance, ImageFilter, ImageFont
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt
from pptx.dml.color import RGBColor

import tkinter as tk
from tkinter import ttk, filedialog, messagebox, scrolledtext

_WIN_TESSERACT = r"C:\Program Files\Tesseract-OCR\tesseract.exe"
pytesseract.pytesseract.tesseract_cmd = os.environ.get(
    "TESSERACT_CMD",
    _WIN_TESSERACT if os.path.exists(_WIN_TESSERACT) else "tesseract",
)

# ════════════════════════════════════════════════════════════════════════
# GUI <-> worker-thread communication
# ════════════════════════════════════════════════════════════════════════

LOG_QUEUE = queue.Queue()
PROGRESS_QUEUE = queue.Queue()
CHOICE_REQUEST_QUEUE = queue.Queue()
CHOICE_RESPONSE_QUEUE = queue.Queue()
RESULT_QUEUE = queue.Queue()


def log(msg=""):
    text = str(msg)
    print(text)
    LOG_QUEUE.put(text)


def set_progress(done, total):
    PROGRESS_QUEUE.put((done, total))


IMAGE_WORKERS = int(os.environ.get("IMAGE_WORKERS", str(min(8, (os.cpu_count() or 4)))))
IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".tiff"}

HIGHLIGHT_COLOR = (237, 108, 66, 40)
HIGHLIGHT_BORDER = (200, 60, 40, 220)
HIGHLIGHT_ACCENT = (200, 45, 35, 255)
OCR_SCALE = 3
OCR_CONTRAST = 2.0
MAX_RESULTS_PER_SHOT = 60
ROW_GAP_BAND = 6
ROW_GAP_SEARCH = 0.4

FLAGGED_KEYWORDS = [
    "court", "courts", "lawsuit", "trial", "judge", "verdict", "conviction",
    "convicted", "sentenced", "sentencing", "plea", "guilty", "acquitted",
    "subpoena", "hearing", "defendant", "plaintiff", "litigation", "docket",
    "indicted", "indictment", "charged", "charges", "arraignment", "bail",
    "parole", "probation", "warrant", "injunction", "restraining order",
    "supreme court", "district court", "federal court", "appeals court",
    "money laundering", "laundering", "wire fraud", "embezzlement",
    "financial crime", "illicit funds", "shell company", "tax evasion",
    "proceeds of crime", "structuring", "smurfing",
    "fraud", "scam", "theft", "robbery", "burglary", "murder", "homicide",
    "assault", "battery", "drug", "narcotics", "trafficking", "smuggling",
    "extortion", "bribery", "corruption", "forgery", "counterfeit",
    "illegal", "criminal", "felony", "misdemeanor", "arrested", "arrest",
    "prison", "jail", "incarcerated", "fugitive", "suspect", "accused",
    "crime", "offense", "violation", "racket", "racketeering", "cartel",
    "gang", "syndicate", "ponzi", "insider trading", "securities fraud",
    "identity theft", "cybercrime", "hacking", "blackmail",
    "transnet",
]


def load_font(size: int) -> ImageFont.FreeTypeFont:
    for path in [
        "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
        "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
        "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
        "/System/Library/Fonts/Helvetica.ttc",
        "C:/Windows/Fonts/arialbd.ttf",
    ]:
        if os.path.exists(path):
            return ImageFont.truetype(path, size)
    return ImageFont.load_default()


def strip_pg_suffix(text: str) -> str:
    """Whatever comes at or after a standalone 'pg' token gets dropped —
    a simple, robust rule that keeps a display name clean regardless of
    exact filename format ('pg 1', 'pg_1', 'PG#1', ...). Applied as a
    final safety net wherever the 'Search' field gets built, on top of
    already using the clean matched keyword rather than raw filename
    text — belt and suspenders, so nothing starting with 'pg' can ever
    leak into the slide."""
    m = re.search(r'\bpg\b', text, re.IGNORECASE)
    if m:
        return text[:m.start()].strip()
    return text.strip()


def extract_name(stem: str) -> str:
    """'Ashley Rocharam pg 1 no 10' -> 'Ashley Rocharam'. Handles both
    space- and underscore-separated filenames ('... pg 1 ...' as well as
    '..._pg_1_...') — matching on '_pg_' alone silently failed to strip
    anything from a space-separated filename, which is this client's
    actual naming convention, leaving the raw 'pg N no M' tail stuck onto
    the end. This function itself isn't used for slide metadata anymore
    (that now always comes from the already-matched clean keyword), but
    kept correct here too as a general-purpose fallback."""
    parts = re.split(r'[\s_]pg[\s_]', stem, maxsplit=1, flags=re.IGNORECASE)
    return parts[0].replace("_", " ").strip()


def extract_page_hint(stem: str) -> int | None:
    """'Sudesh Premchand Rocharam pg 1 no 1-2,5,8' -> 1. Accepts a space,
    '#', or underscore between 'pg' and the number ('pg_36' as well as
    'pg 36'/'pg#36') — an underscore-separated filename used to silently
    fail this before, breaking the page_hint fallback for any screenshot
    whose page-number text Bing doesn't print at all."""
    m = re.search(r'pg\s*[#_]?\s*(\d+)', stem, re.IGNORECASE)
    return int(m.group(1)) if m else None


def extract_flagged_rows_from_filename(stem: str) -> list[int] | None:
    """'Ashley Rocharam pg 5 no 2-4,8' -> [2, 3, 4, 8]. This is now the
    SOLE source of which rows get flagged — the person naming the file
    has already reviewed the screenshot and is directly recording which
    results matter, which is more reliable than an automated keyword
    scan (no OCR run needed here at all, which is also the slowest step
    in the whole pipeline — a real speed win alongside being correct).
    Ranges expand ('2-4' -> 2,3,4); comma or ampersand both work as
    separators ('no 6,8&10'), matching every naming convention seen so
    far. Returns None if no 'no' clause is found at all — an empty
    result and a missing clause are kept distinct (empty list still
    marks the screenshot as reviewed-and-clean, not undetected)."""
    m = re.search(r'\bno\s*[#]?\s*([\d,&\-\s]+)', stem, re.IGNORECASE)
    if not m:
        return None
    rows = set()
    for part in re.split(r'[,&]', m.group(1)):
        part = part.strip()
        if not part:
            continue
        range_m = re.match(r'^(\d+)\s*-\s*(\d+)', part)
        if range_m:
            lo, hi = int(range_m.group(1)), int(range_m.group(2))
            if lo <= hi:
                rows.update(range(lo, hi + 1))
            continue
        digit_m = re.match(r'^(\d+)', part)
        if digit_m:
            rows.add(int(digit_m.group(1)))
    return sorted(rows) if rows else None


def row_is_flagged(ocr_text: str) -> bool:
    text_lower = ocr_text.lower()
    return any(kw in text_lower for kw in FLAGGED_KEYWORDS)


def ocr_row(img: Image.Image, y0: int, y1: int, x_right: int | None = None) -> str:
    crop = img.crop((160, y0, x_right or img.width, y1))
    gray = crop.convert("L")
    enhanced = ImageEnhance.Contrast(gray).enhance(OCR_CONTRAST)
    sharpened = enhanced.filter(ImageFilter.SHARPEN)
    big = sharpened.resize(
        (sharpened.width * OCR_SCALE, sharpened.height * OCR_SCALE), Image.LANCZOS
    )
    return pytesseract.image_to_string(big, config="--psm 6 --oem 3")


RESULT_RANGE_PATTERN = re.compile(
    r'(\d[\d,]*)\s*[-–—]\s*(\d[\d,]*)\s+[o0]f\s+[\d,]+', re.IGNORECASE
)
RESULT_COUNT_PATTERN = re.compile(r'\d[\d,]*\s+results', re.IGNORECASE)


def detect_page_info(img: Image.Image, page_hint: int | None = None,
                      results_per_page: int = 10) -> tuple[int, int, int]:
    H = img.height
    scan_to = min(250, int(H * 0.20))
    crop = img.crop((80, 60, 700, scan_to + 40))
    big = crop.resize((crop.width * 2, crop.height * 2), Image.LANCZOS)
    text = pytesseract.image_to_string(big, config="--psm 6")

    for m in RESULT_RANGE_PATTERN.finditer(text):
        try:
            r_start = int(m.group(1).replace(",", ""))
            r_end = int(m.group(2).replace(",", ""))
        except ValueError:
            continue
        total = r_end - r_start + 1
        if 1 <= total <= MAX_RESULTS_PER_SHOT:
            page = (r_start - 1) // total + 1
            return page, r_start, r_end

    if page_hint is not None:
        r_start = (page_hint - 1) * results_per_page + 1
        r_end = page_hint * results_per_page
        return page_hint, r_start, r_end

    if RESULT_COUNT_PATTERN.search(text):
        return 1, 1, 10

    log("  ⚠  Could not reliably detect result range — defaulting to Page 1, 1–10")
    return 1, 1, 10


def detect_results_area(arr: np.ndarray) -> tuple[int, int]:
    H = arr.shape[0]
    scan_start = int(H * 0.08)
    scan_end = int(H * 0.95)
    R = arr[scan_start:scan_end, 160:560, 0].astype(int)
    G = arr[scan_start:scan_end, 160:560, 1].astype(int)
    B = arr[scan_start:scan_end, 160:560, 2].astype(int)
    blue_rows = np.where(np.sum((R < 80) & (G < 140) & (B > 130), axis=1) > 5)[0]
    y_top = (int(blue_rows[0]) + scan_start - 20) if len(blue_rows) else scan_start
    return max(0, y_top), int(H * 0.93)


def refine_first_row_top(arr: np.ndarray, y_top_estimate: int, x_right: int, max_pad: int = 90) -> int:
    darkness = compute_darkness_profile(arr, x_right)
    floor = max(0, y_top_estimate - max_pad - 40)
    return anchor_to_boundary(darkness, y_top_estimate, floor, max_pad)


def content_right_edge(arr: np.ndarray, y_top: int, y_bottom: int, W: int) -> int:
    x0 = int(W * 0.45)
    col = arr[y_top:y_bottom, x0:W]
    nonwhite = np.any(col < 225, axis=2).sum(axis=0)
    white = nonwhite < 0.02 * max(1, (y_bottom - y_top))
    min_gutter = int(0.035 * W)
    run = 0
    start = 0
    for i, w in enumerate(white):
        if w:
            if run == 0:
                start = i
            run += 1
            if run >= min_gutter:
                return x0 + start
        else:
            run = 0
    return W


def trim_bottom(arr: np.ndarray, y_top: int, y_bottom: int, x_right: int) -> int:
    """
    Pull y_bottom up past trailing whitespace so the last row doesn't
    stretch down into the footer/pagination area.
    """
    col = arr[y_top:y_bottom, 160:x_right]
    ink = np.any(col < 210, axis=2).sum(axis=1) > 3
    rows = np.where(ink)[0]
    if len(rows) == 0:
        return y_bottom

    last_ink = int(rows[-1])
    search_from = rows[0] + int((last_ink - rows[0]) * 0.5)

    best_gap = None
    gap_start = None
    for i in range(search_from, last_ink + 1):
        if not ink[i]:
            if gap_start is None:
                gap_start = i
        else:
            if gap_start is not None:
                gap_len = i - gap_start
                if gap_len >= 18 and (best_gap is None or gap_len > best_gap[1]):
                    best_gap = (gap_start, gap_len)
                gap_start = None

    if best_gap is not None:
        return y_top + best_gap[0] + 12
    return y_top + last_ink + 12


def compute_darkness_profile(arr: np.ndarray, x_right: int) -> np.ndarray:
    x0 = 160
    x1 = max(x0 + 1, min(x_right, arr.shape[1]))
    gray = arr[:, x0:x1, :].astype(int).mean(axis=2)
    return (255 - gray).sum(axis=1)


def snap_row_boundaries(arr: np.ndarray, boundaries: list[int], x_right: int) -> list[int]:
    if len(boundaries) < 3:
        return boundaries
    darkness = compute_darkness_profile(arr, x_right)
    step = boundaries[1] - boundaries[0]
    radius = max(ROW_GAP_BAND, int(step * ROW_GAP_SEARCH))
    snapped = boundaries[:]
    for i in range(1, len(boundaries) - 1):
        est = boundaries[i]
        lo = max(0, boundaries[i - 1] + ROW_GAP_BAND, est - radius)
        hi = min(len(darkness), boundaries[i + 1] - ROW_GAP_BAND, est + radius)
        if hi - lo <= ROW_GAP_BAND:
            continue
        window = darkness[lo:hi]
        smoothed = np.convolve(window, np.ones(ROW_GAP_BAND), mode="valid")
        best = int(np.argmin(smoothed)) + ROW_GAP_BAND // 2
        snapped[i] = lo + best
    return snapped


def anchor_to_boundary(darkness: np.ndarray, anchor_y: int, floor: int, max_pad: int) -> int:
    lo = max(floor, anchor_y - max_pad)
    hi = anchor_y
    if hi - lo <= ROW_GAP_BAND:
        return min(hi, max(floor, anchor_y - max_pad))

    window = darkness[lo:hi]
    threshold = max(float(np.percentile(window, 25)), 1.0)
    blank = window <= threshold

    run_start = None
    run_len = 0
    best = None
    for i, is_blank in enumerate(blank):
        if is_blank:
            if run_start is None:
                run_start = i
            run_len += 1
            if run_len >= ROW_GAP_BAND and best is None:
                best = run_start + run_len // 2
        else:
            run_start = None
            run_len = 0

    if best is not None:
        return lo + best

    smoothed = np.convolve(window, np.ones(ROW_GAP_BAND), mode="valid")
    best_idx = int(np.argmin(smoothed)) + ROW_GAP_BAND // 2
    return lo + best_idx


def _merge_undersized_rows(boundaries: list[int], min_height_ratio: float = 0.4) -> list[int]:
    if len(boundaries) <= 2:
        return boundaries
    changed = True
    while changed and len(boundaries) > 2:
        changed = False
        heights = [boundaries[i + 1] - boundaries[i] for i in range(len(boundaries) - 1)]
        avg = sum(heights) / len(heights)
        for i, h in enumerate(heights):
            if h < avg * min_height_ratio:
                if i == 0:
                    del boundaries[1]
                else:
                    del boundaries[i]
                changed = True
                break
    return boundaries


def detect_row_boundaries_by_gaps(arr: np.ndarray, y_top: int, y_bottom: int,
                                   x_right: int, min_gap: int = 28) -> list[int] | None:
    col = arr[y_top:y_bottom, 160:x_right]
    ink = np.any(col < 210, axis=2).sum(axis=1) > 3
    n = len(ink)

    gaps = []
    gap_start = None
    for i in range(n):
        if not ink[i]:
            if gap_start is None:
                gap_start = i
        else:
            if gap_start is not None:
                gaps.append((gap_start, i))
                gap_start = None
    if gap_start is not None:
        gaps.append((gap_start, n))

    interior = []
    for g0, g1 in gaps:
        if g1 - g0 < min_gap:
            continue
        if g0 <= 2 or g1 >= n - 2:
            continue
        interior.append(y_top + (g0 + g1) // 2)

    if not interior:
        return None
    boundaries = [y_top] + interior + [y_bottom]
    return _merge_undersized_rows(boundaries)


def detect_row_boundaries(arr: np.ndarray, y_top: int, y_bottom: int, total: int, W: int,
                           x_right: int) -> list[int]:
    avg = (y_bottom - y_top) / total
    even = [int(y_top + i * avg) for i in range(total + 1)]

    x0, x1 = 160, int(W * 0.50)
    sub = arr[y_top:y_bottom, x0:x1]
    R, G, B = sub[..., 0].astype(int), sub[..., 1].astype(int), sub[..., 2].astype(int)
    blue = np.sum((R < 90) & (G < 150) & (B > 120), axis=1) > 4
    rows = np.where(blue)[0]
    if len(rows) < 2:
        return snap_row_boundaries(arr, even, x_right)

    anchors = [int(rows[0])]
    for a, b in zip(rows[:-1], rows[1:]):
        if b - a > 25:
            anchors.append(int(b))
    merged = [anchors[0]]
    for a in anchors[1:]:
        if a - merged[-1] >= 0.45 * avg:
            merged.append(a)
    anchors = merged
    while len(anchors) > total:
        j = int(np.argmin(np.diff(anchors)))
        del anchors[j + 1]
    if len(anchors) != total:
        return snap_row_boundaries(arr, even, x_right)

    max_pad = int(min(70, max(45, 0.5 * avg)))
    darkness = compute_darkness_profile(arr, x_right)
    bounds = []
    floor = y_top
    for a in anchors:
        b = anchor_to_boundary(darkness, y_top + a, floor, max_pad)
        bounds.append(b)
        floor = b
    bounds.append(y_bottom)
    return bounds


def ink_extent(arr: np.ndarray, top: int, bot: int, x_right: int) -> tuple[int, int] | None:
    reg = arr[top:bot, 60:x_right]
    inked = np.any(reg < 210, axis=2).sum(axis=1) > 3
    idx = np.where(inked)[0]
    if len(idx) == 0:
        return None
    gap_limit = 55
    start = int(idx[0])
    end = start
    white = 0
    for y in range(start, len(inked)):
        if inked[y]:
            end = y
            white = 0
        else:
            white += 1
            if white > gap_limit:
                break
    if end - start < 18:
        return None
    return top + start - 4, top + end + 4


def make_badge(result_num: int, font) -> Image.Image:
    label = f"#{result_num}"
    bbox = font.getbbox(label)
    tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
    px, py = 14, 8
    bw, bh = tw + px * 2, th + py * 2
    badge = Image.new("RGBA", (bw, bh), (0, 0, 0, 0))
    bd = ImageDraw.Draw(badge)
    bd.rounded_rectangle([0, 0, bw - 1, bh - 1], radius=bh // 2, fill=(25, 65, 145, 220))
    bd.text((px, py), label, font=font, fill=(255, 255, 255, 255))
    return badge


def make_header(width: int, page: int, r_start: int, r_end: int, part: int, of_parts: int) -> Image.Image:
    h = 60
    header = Image.new("RGB", (width, h), (25, 65, 145))
    draw = ImageDraw.Draw(header)
    draw.text((18, 18), f"Page {page}", font=load_font(22), fill=(255, 255, 255))
    draw.text((width // 2, 20), f"Results {r_start}–{r_end}", font=load_font(15), fill=(180, 215, 255), anchor="mt")
    draw.text((width - 18, 18), f"Part {part} of {of_parts}", font=load_font(22), fill=(255, 255, 255), anchor="ra")
    return header


def apply_highlight(canvas: Image.Image, x0: int, y0: int, x1: int, y1: int) -> None:
    """A cleaner, more professional flagged-row treatment than a flat
    highlighter-yellow box: a muted warm-red tint (reads as "flagged for
    review" rather than "highlighter pen"), a full rounded border on all
    four sides instead of just top/bottom rules, and a solid accent bar
    down the left edge — the same visual language document-review tools
    use, so a flagged row is unambiguous even skimmed quickly or printed
    in black and white (the accent bar and border still show as a solid
    dark shape even with color stripped)."""
    overlay = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    d = ImageDraw.Draw(overlay)
    d.rounded_rectangle([x0, y0, x1, y1], radius=6, fill=HIGHLIGHT_COLOR,
                         outline=HIGHLIGHT_BORDER, width=2)
    accent_w = 6
    d.rounded_rectangle([x0, y0, x0 + accent_w, y1], radius=3, fill=HIGHLIGHT_ACCENT)
    canvas.alpha_composite(overlay)


# ════════════════════════════════════════════════════════════════════════
# Detection-only pass (review-before-build) + full build
# ════════════════════════════════════════════════════════════════════════

def compute_auto_flagged(img: Image.Image, boundaries: list[int], x_right: int, total: int) -> set[int]:
    auto_flagged = set()
    for i in range(total):
        text = ocr_row(img, boundaries[i], boundaries[i + 1], x_right)
        if row_is_flagged(text):
            auto_flagged.add(i + 1)
    return auto_flagged


def compute_auto_flagged_from_spans(img: Image.Image, row_spans: list[tuple[int, int]], x_right: int) -> set[int]:
    """Same as compute_auto_flagged, but takes each row's own (top,
    bottom) span directly rather than a shared boundaries list — used
    on the precomputed_rows path, where rows may no longer be
    contiguous (a deleted row leaves a gap), so there's no valid single
    boundaries list to index into in the first place."""
    auto_flagged = set()
    for i, (top, bot) in enumerate(row_spans):
        text = ocr_row(img, top, bot, x_right)
        if row_is_flagged(text):
            auto_flagged.add(i + 1)
    return auto_flagged


def detect_rows_for_review(image_path, page_hint: int | None = None) -> dict:
    image_path = Path(image_path)
    img = Image.open(image_path).convert("RGB")
    W, H = img.size
    arr = np.array(img)

    if page_hint is None:
        page_hint = extract_page_hint(image_path.stem)
    page, r_first, r_last = detect_page_info(img, page_hint=page_hint)
    ocr_total = r_last - r_first + 1
    if ocr_total < 2:
        log(f"  ⚠  {image_path.name}: detected range too small ({r_first}-{r_last}), defaulting to 1-10")
        page, r_first, r_last = 1, 1, 10
        ocr_total = 10

    y_top, y_bottom = detect_results_area(arr)
    x_right = content_right_edge(arr, y_top, y_bottom, W)
    y_bottom = trim_bottom(arr, y_top, y_bottom, x_right)
    y_top = refine_first_row_top(arr, y_top, x_right)

    gap_boundaries = detect_row_boundaries_by_gaps(arr, y_top, y_bottom, x_right)
    if gap_boundaries is not None:
        boundaries = gap_boundaries
        total = len(boundaries) - 1
        if total != ocr_total:
            log(f"  ℹ  {image_path.name}: whitespace-gap detection found {total} row(s) — "
                f"header text suggested {ocr_total}, using the {total} actually detected.")
        r_last = r_first + total - 1
    else:
        total = ocr_total
        boundaries = detect_row_boundaries(arr, y_top, y_bottom, total, W, x_right)

    row_coords = [{"row": i + 1, "top": int(boundaries[i]), "bottom": int(boundaries[i + 1])} for i in range(total)]
    filename_flagged = extract_flagged_rows_from_filename(image_path.stem)
    # Clamp to the valid row range — a human typed this by hand, and a
    # typo (a row number past what's actually on the page) shouldn't
    # silently break anything downstream; it just can't apply.
    auto_flagged = {n for n in (filename_flagged or []) if 1 <= n <= total}
    if filename_flagged and any(n > total for n in filename_flagged):
        log(f"  ⚠  {image_path.name}: filename lists row(s) beyond the {total} detected "
            f"— {[n for n in filename_flagged if n > total]} ignored.")

    return {
        "total_rows": total, "x_right": x_right, "page": page,
        "r_first": r_first, "r_last": r_last, "row_coords": row_coords,
        "suggested_flagged": sorted(auto_flagged), "name": extract_name(image_path.stem),
    }


def process_explorer_image(image_path, out_dir, page: int | None = None,
                            r_first: int | None = None, r_last: int | None = None,
                            confirmed_flagged: list[int] | None = None,
                            precomputed_rows: list[tuple[int, int]] | None = None,
                            precomputed_x_right: int | None = None) -> dict:
    image_path = Path(image_path)
    img = Image.open(image_path).convert("RGB")
    W, H = img.size
    arr = np.array(img)

    if precomputed_rows is not None and precomputed_x_right is not None:
        if page is None or r_first is None or r_last is None:
            raise ValueError("page/r_first/r_last are required together with precomputed_rows")
        # Each row's OWN (top, bottom) is used directly — never rebuilt
        # into one shared "boundaries" list. A shared list only works
        # when every row's bottom edge exactly touches the next row's
        # top edge, which is true for a fresh, untouched detection but
        # becomes FALSE the moment a row is deleted in review_report_
        # explorer.py: the two rows now on either side of the gap are no
        # longer adjacent, so reconstructing "boundaries" as
        # [first.top] + [row.bottom for row in rows] silently stitches
        # across that gap — verified on a real screenshot: deleting the
        # "IDAC Report" row still left its exact content sitting in the
        # output, merged into the next kept row's box. Keeping explicit
        # (top, bottom) pairs per row sidesteps the whole issue: there is
        # no shared list to go stale.
        row_spans = [(int(t), int(b)) for t, b in precomputed_rows]
        x_right = precomputed_x_right
        total = len(row_spans)
    else:
        page_hint = extract_page_hint(image_path.stem)
        page, r_first, r_last = detect_page_info(img, page_hint=page_hint)
        ocr_total = r_last - r_first + 1
        if ocr_total < 2:
            page, r_first, r_last = 1, 1, 10
            ocr_total = 10
        y_top, y_bottom = detect_results_area(arr)
        x_right = content_right_edge(arr, y_top, y_bottom, W)
        y_bottom = trim_bottom(arr, y_top, y_bottom, x_right)
        y_top = refine_first_row_top(arr, y_top, x_right)
        gap_boundaries = detect_row_boundaries_by_gaps(arr, y_top, y_bottom, x_right)
        if gap_boundaries is not None:
            boundaries = gap_boundaries
            total = len(boundaries) - 1
            r_last = r_first + total - 1
        else:
            total = ocr_total
            boundaries = detect_row_boundaries(arr, y_top, y_bottom, total, W, x_right)
        # a fresh detection IS contiguous by construction, so converting
        # the boundaries list into explicit spans here is always safe
        row_spans = [(boundaries[i], boundaries[i + 1]) for i in range(total)]

    if confirmed_flagged is not None:
        final_flagged_set = {n for n in confirmed_flagged if 1 <= n <= total}
    else:
        filename_flagged = extract_flagged_rows_from_filename(image_path.stem)
        final_flagged_set = {n for n in (filename_flagged or []) if 1 <= n <= total}

    name = extract_name(image_path.stem)
    badge_font = load_font(20)
    label_font = load_font(17)
    size_part1 = (total + 1) // 2
    part_sizes = [size_part1, total - size_part1]
    stem = image_path.stem
    outputs = []
    highlighted = []
    row_coords_out = []

    r_off = 0
    for part_idx, part_size in enumerate(part_sizes):
        if part_size <= 0:
            continue
        part_num = part_idx + 1
        header_h = 60

        # Crop EACH row's own tight slice individually (using its own
        # span, never a shared/rebuilt boundaries list) and stack them
        # directly together — this is what makes a deleted row's content
        # genuinely absent from the output, not just unlabeled. Cropped
        # to x_right (the real results column), NOT the full page width
        # — a page with a side panel (verified on a real Bing screenshot:
        # an "AI Answer" box running down the right side) has that panel
        # spanning a much taller continuous region than any single row,
        # so slicing per-row at full width and re-stacking chops it up
        # into visibly scrambled fragments. The results column itself
        # doesn't need that panel anyway — it isn't a real search result.
        part_spans = row_spans[r_off:r_off + part_size]
        row_slices = [img.crop((0, top, x_right, bot)) for top, bot in part_spans]
        total_height = sum(rs.height for rs in row_slices)

        canvas = Image.new("RGBA", (x_right, header_h + total_height), (255, 255, 255, 255))
        canvas.paste(make_header(x_right, page, r_first + r_off, r_first + r_off + part_size - 1,
                                  part_num, len(part_sizes)), (0, 0))

        y_cursor = header_h
        for i, row_slice in enumerate(row_slices):
            canvas.paste(row_slice.convert("RGBA"), (0, y_cursor))
            row_top, row_bot = part_spans[i]
            local_num = r_off + i + 1
            local_top = y_cursor
            local_bot = y_cursor + row_slice.height
            row_coords_out.append({"row": local_num, "part": part_num,
                                    "top": int(row_top), "bottom": int(row_bot)})

            if local_num in final_flagged_set:
                ext = ink_extent(arr, row_top, row_bot, x_right)
                if ext is not None:
                    # ext is measured in the ORIGINAL image's coordinates
                    # (ink_extent scans the full raw arr) — translate it
                    # into this row's own local position on the canvas,
                    # since rows no longer sit at their original absolute
                    # Y position once gaps have been squeezed out.
                    hy0 = local_top + (ext[0] - row_top)
                    hy1 = local_top + (ext[1] - row_top)
                else:
                    hy0, hy1 = local_top, local_bot
                apply_highlight(canvas, 0, hy0, x_right, hy1)
                highlighted.append(local_num)

            if i > 0:
                draw = ImageDraw.Draw(canvas)
                draw.line([(160, local_top), (x_right - 20, local_top)], fill=(200, 200, 200, 200), width=1)

            canvas.alpha_composite(make_badge(local_num, badge_font), (20, local_top + 12))
            draw = ImageDraw.Draw(canvas)
            draw.text((x_right - 22, local_top + 10), f"Result {local_num} of {total}", font=label_font,
                      fill=(100, 100, 100, 220), anchor="ra")

            y_cursor += row_slice.height

        out_path = Path(out_dir) / f"{stem}_part{part_num}.jpg"
        canvas.convert("RGB").save(out_path, "JPEG", quality=92)
        outputs.append({"path": str(out_path), "part": part_num, "of": len(part_sizes),
                         "r_start": r_first + r_off, "r_end": r_first + r_off + part_size - 1})
        r_off += part_size

    return {
        "file": image_path.name, "name": name, "page": page, "r_first": r_first, "r_last": r_last,
        "total_rows": total, "x_right": x_right, "highlighted": sorted(highlighted),
        "outputs": outputs, "row_coords": row_coords_out, "raw_filepath": str(image_path),
    }


# ════════════════════════════════════════════════════════════════════════
# Folder / filesystem helpers
# ════════════════════════════════════════════════════════════════════════

MONTH_NAMES = [
    "January", "February", "March", "April", "May", "June",
    "July", "August", "September", "October", "November", "December",
]
MAIN_FOLDER_NAME_HINT = "Sudesh"


def list_dirs(path):
    try:
        entries = os.listdir(path)
    except FileNotFoundError:
        return []
    return sorted(e for e in entries if os.path.isdir(os.path.join(path, e)))


def list_files(path):
    try:
        entries = os.listdir(path)
    except FileNotFoundError:
        return []
    return sorted(e for e in entries if os.path.isfile(os.path.join(path, e)))


def prompt_choice(question, options):
    if not options:
        return None
    if len(options) == 1:
        return options[0]
    CHOICE_REQUEST_QUEUE.put((question, list(options)))
    return CHOICE_RESPONSE_QUEUE.get()


def extract_day_number(folder_name):
    match = re.match(r'^\s*(\d{1,2})', folder_name.strip())
    if match:
        day = int(match.group(1))
        if 1 <= day <= 31:
            return day
    return None


def week_for_day(day_num):
    if 1 <= day_num <= 7:
        return 1
    elif 8 <= day_num <= 14:
        return 2
    elif 15 <= day_num <= 21:
        return 3
    elif day_num >= 22:
        return 4
    return None


def find_subfolder(path, name_hint):
    for d in list_dirs(path):
        if d.strip().lower() == name_hint.lower():
            return d
    return None


# ════════════════════════════════════════════════════════════════════════
# Keyword matching / CSV
# ════════════════════════════════════════════════════════════════════════

KEYWORDS = [
    "Prashin Sudesh Premchand Rocharam",
    "Prashin Premchand Rocharam",
    "Sudesh Premchand Rocharam",
    "Sudesh Rocharam",
    "Ashley Premchand Rocharam",
    "Ashley Rocharam",
]
KEYWORDS_SORTED = sorted(KEYWORDS, key=len, reverse=True)

NO_RESULTS_PATTERN = re.compile(r'no[\s_]*results', re.IGNORECASE)
MARKED_PATTERN = re.compile(r'marked', re.IGNORECASE)


def is_no_results(filename_no_ext: str) -> bool:
    return bool(NO_RESULTS_PATTERN.search(filename_no_ext))


def is_marked(filename_no_ext: str) -> bool:
    return bool(MARKED_PATTERN.search(filename_no_ext))


ROW_ORDER = [
    "Google/Explorer",
    "A. Prashin Sudesh Premchand Rocharam",
    "B. Prashin Premchand Rocharam",
    "C. Sudesh Premchand Rocharam",
    "D. Sudesh Rocharam",
    "E. Ashley Premchand Rocharam",
    "F. Ashley Rocharam",
]
ROW_LABEL_TO_KEYWORD = {
    "A. Prashin Sudesh Premchand Rocharam": "Prashin Sudesh Premchand Rocharam",
    "B. Prashin Premchand Rocharam": "Prashin Premchand Rocharam",
    "C. Sudesh Premchand Rocharam": "Sudesh Premchand Rocharam",
    "D. Sudesh Rocharam": "Sudesh Rocharam",
    "E. Ashley Premchand Rocharam": "Ashley Premchand Rocharam",
    "F. Ashley Rocharam": "Ashley Rocharam",
}
HEADER = ["Tags - Keywords", "Week 1 dd-mm-yy", "Week 2 dd-mm-yy", "Week 3 dd-mm-yy", "Week 4 dd-mm-yy"]


def match_keyword(text):
    for kw in KEYWORDS_SORTED:
        if kw.lower() in text.lower():
            return kw
    return None


def pgno_label(page: int, numbers: list[int]) -> str:
    if not numbers:
        return f"Pg{page}"
    return f"Pg{page}#{','.join(str(n) for n in numbers)}"


def build_csv_rows(csv_results):
    rows = [HEADER]
    for label in ROW_ORDER:
        keyword = ROW_LABEL_TO_KEYWORD.get(label)
        row = [label]
        for week_num in [1, 2, 3, 4]:
            if keyword:
                entries = csv_results.get(keyword, {}).get(week_num, [])
                row.append(", ".join(entries))
            else:
                row.append("")
        rows.append(row)
    return rows


def write_csv_report(report_root, csv_results, filename="main file.csv"):
    os.makedirs(report_root, exist_ok=True)
    report_path = os.path.join(report_root, filename)
    with open(report_path, 'w', newline='', encoding='utf-8') as f:
        csv.writer(f).writerows(build_csv_rows(csv_results))
    return report_path


# ════════════════════════════════════════════════════════════════════════
# PPTX generation
# ════════════════════════════════════════════════════════════════════════

IDX_TITLE = 0
IDX_DIVIDER = 2
IDX_RESULT = 3
IDX_CONCLUSION = 100

CONTENT_TOP = 457200
CONTENT_HEIGHT = 5923584
RIGHT_MARGIN = 0


def duplicate_slide(prs, index):
    source = prs.slides[index]
    dest = prs.slides.add_slide(source.slide_layout)
    for shp in list(dest.shapes):
        shp._element.getparent().remove(shp._element)

    rid_map = {}
    for rel_id, rel in source.part.rels.items():
        if rel.is_external:
            continue
        rid_map[rel_id] = dest.part.relate_to(rel.target_part, rel.reltype)

    for shape_el in list(source.shapes._spTree):
        if shape_el.tag in (qn('p:nvGrpSpPr'), qn('p:grpSpPr')):
            continue
        new_el = copy.deepcopy(shape_el)
        for el in new_el.iter():
            for attr in ('embed', 'link', 'id'):
                full_attr = qn('r:' + attr)
                if full_attr in el.attrib and el.attrib[full_attr] in rid_map:
                    el.attrib[full_attr] = rid_map[el.attrib[full_attr]]
        dest.shapes._spTree.append(new_el)
    return dest


def set_title_text(slide, text):
    for shp in slide.shapes:
        if shp.name.startswith("Title") and shp.has_text_frame:
            tf = shp.text_frame
            first_p = tf.paragraphs[0]
            if first_p.runs:
                first_p.runs[0].text = text
                for extra in first_p.runs[1:]:
                    extra.text = ""
            else:
                first_p.text = text
            for p in list(tf.paragraphs[1:]):
                p._p.getparent().remove(p._p)
            return True
    return False


def set_body_text(slide, text, placeholder_name_prefix="Content Placeholder"):
    for shp in slide.shapes:
        if shp.name.startswith(placeholder_name_prefix) and shp.has_text_frame:
            tf = shp.text_frame
            first_p = tf.paragraphs[0]
            if first_p.runs:
                first_p.runs[0].text = text
                for extra in first_p.runs[1:]:
                    extra.text = ""
            else:
                first_p.text = text
            for p in list(tf.paragraphs[1:]):
                p._p.getparent().remove(p._p)
            return True
    return False


def set_metadata_text(slide, fields):
    for shp in slide.shapes:
        if shp.name.startswith("Title") and shp.has_text_frame:
            tf = shp.text_frame
            tf.clear()
            p = tf.paragraphs[0]
            for i, (label, value) in enumerate(fields):
                if i > 0:
                    p.add_line_break()
                r_label = p.add_run()
                r_label.text = f"{label}: "
                r_label.font.bold = True
                r_label.font.size = Pt(16)
                r_label.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                r_value = p.add_run()
                r_value.text = str(value)
                r_value.font.bold = False
                r_value.font.size = Pt(16)
                r_value.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
            return True
    return False


def delete_slide(prs, index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    rId = slides[index].get(qn('r:id'))
    prs.part.drop_rel(rId)
    xml_slides.remove(slides[index])


def replace_picture(slide, image_path, slide_width):
    for shp in list(slide.shapes):
        if shp.shape_type == 13:
            shp._element.getparent().remove(shp._element)

    zone_left = slide_width // 3
    zone_width = slide_width - zone_left

    with Image.open(image_path) as im:
        iw, ih = im.size
    aspect = iw / ih

    height = CONTENT_HEIGHT
    width = int(height * aspect)
    if width > zone_width:
        width = zone_width
        height = int(width / aspect)

    left = slide_width - width - RIGHT_MARGIN
    top = CONTENT_TOP + (CONTENT_HEIGHT - height) // 2
    slide.shapes.add_picture(image_path, left, top, width=width, height=height)


def build_pptx(template_path, out_path, month_name, year, week_groups, summary_text):
    prs = Presentation(template_path)
    n_original = len(prs.slides._sldIdLst)
    slide_w = prs.slide_width

    title_slide = duplicate_slide(prs, IDX_TITLE)
    set_title_text(title_slide, "SEO REPORT")
    for shp in title_slide.shapes:
        if shp.name.startswith("Subtitle") and shp.has_text_frame:
            shp.text_frame.paragraphs[0].runs[0].text = f"{month_name} {year} — Explorer"

    for week_num in sorted(week_groups):
        entries = week_groups[week_num]
        if not entries:
            continue
        week_divider = duplicate_slide(prs, IDX_DIVIDER)
        set_title_text(week_divider, f"WEEK {week_num} SEARCH RESULTS")

        for entry in entries:
            rs = duplicate_slide(prs, IDX_RESULT)
            replace_picture(rs, entry["path"], slide_w)
            page_item = entry["pgno"] or f"Results {entry.get('r_start', '')}–{entry.get('r_end', '')}"
            set_metadata_text(rs, [
                ("Date", entry["date"]),
                ("Page/Item No.", page_item),
                ("Search", entry["keyword"]),
                ("Part", f"{entry['part']} of {entry['of']}"),
            ])

    concl = duplicate_slide(prs, IDX_CONCLUSION)
    set_body_text(concl, summary_text)

    for i in range(n_original - 1, -1, -1):
        delete_slide(prs, i)

    prs.save(out_path)
    return out_path


# ════════════════════════════════════════════════════════════════════════
# Manifest (review-before-build)
# ════════════════════════════════════════════════════════════════════════

def _paths_to_relative(entries, anchor):
    out = []
    for entry in entries:
        e = dict(entry)
        try:
            e["raw_filepath"] = os.path.relpath(entry["raw_filepath"], anchor)
        except ValueError:
            pass
        new_outputs = []
        for o in entry.get("outputs", []):
            oo = dict(o)
            try:
                oo["path"] = os.path.relpath(o["path"], anchor)
            except ValueError:
                pass
            new_outputs.append(oo)
        e["outputs"] = new_outputs
        out.append(e)
    return out


def _paths_to_absolute(entries, anchor):
    out = []
    for entry in entries:
        e = dict(entry)
        if not os.path.isabs(e["raw_filepath"]):
            e["raw_filepath"] = os.path.normpath(os.path.join(anchor, e["raw_filepath"]))
        new_outputs = []
        for o in entry.get("outputs", []):
            oo = dict(o)
            if oo.get("path") and not os.path.isabs(oo["path"]):
                oo["path"] = os.path.normpath(os.path.join(anchor, oo["path"]))
            new_outputs.append(oo)
        e["outputs"] = new_outputs
        out.append(e)
    return out


def save_manifest(report_root, entries, month_name, year, month_num, template_path=None):
    anchor = os.path.dirname(report_root)
    rel_template = template_path
    if template_path:
        try:
            rel_template = os.path.relpath(template_path, anchor)
        except ValueError:
            pass
    data = {
        "month_name": month_name, "year": year, "month_num": month_num,
        "template_path": rel_template, "entries": _paths_to_relative(entries, anchor),
    }
    path = os.path.join(report_root, "review_manifest.json")
    os.makedirs(report_root, exist_ok=True)
    with open(path, "w", encoding="utf-8") as f:
        json.dump(data, f, indent=2)
    return path


def load_manifest(report_root):
    path = os.path.join(report_root, "review_manifest.json")
    if not os.path.exists(path):
        return None
    with open(path, "r", encoding="utf-8") as f:
        data = json.load(f)
    anchor = os.path.dirname(report_root)
    data["entries"] = _paths_to_absolute(data["entries"], anchor)
    if data.get("template_path") and not os.path.isabs(data["template_path"]):
        data["template_path"] = os.path.normpath(os.path.join(anchor, data["template_path"]))
    return data


def build_report_from_manifest(report_root, manifest, template_path=None, extra_csv_results=None):
    images_dir = os.path.join(report_root, "images")
    os.makedirs(images_dir, exist_ok=True)

    csv_results = defaultdict(lambda: defaultdict(list))
    if extra_csv_results:
        for kw, weeks in extra_csv_results.items():
            for wk, vals in weeks.items():
                csv_results[kw][wk].extend(vals)
    week_groups = defaultdict(list)
    total_highlighted = 0
    total_images = 0

    for entry in manifest["entries"]:
        raw_path = Path(entry["raw_filepath"])
        if not raw_path.exists():
            log(f"  ✗  raw screenshot missing, skipping: {raw_path}")
            continue

        for old_out in entry.get("outputs", []):
            old_path = old_out.get("path")
            if old_path and os.path.exists(old_path):
                try:
                    os.remove(old_path)
                except OSError:
                    pass

        precomputed_rows = [(rc["top"], rc["bottom"]) for rc in entry["row_coords"]]
        # r_last needs to reflect the CURRENT row count, not whatever it
        # was at the original detection — a reviewer edit that adds or
        # removes rows changes how many results this screenshot actually
        # covers, and the "Results N-M" label/CSV page reference should
        # track that, not a stale pre-edit number.
        adjusted_r_last = entry["r_first"] + len(precomputed_rows) - 1
        result = process_explorer_image(
            raw_path, images_dir, page=entry["page"], r_first=entry["r_first"], r_last=adjusted_r_last,
            confirmed_flagged=entry["highlighted"],
            precomputed_rows=precomputed_rows, precomputed_x_right=entry["x_right"],
        )
        entry["outputs"] = result["outputs"]
        entry["row_coords"] = result["row_coords"]
        entry["total_rows"] = result["total_rows"]
        entry["r_last"] = result["r_last"]
        entry["highlighted"] = result["highlighted"]
        entry["file"] = result["file"]

        total_images += 1
        total_highlighted += len(entry["highlighted"])
        detected_pgno = pgno_label(entry["page"], entry["highlighted"])
        csv_results[entry["keyword"]][entry["week_num"]].append(detected_pgno)
        for out in entry["outputs"]:
            week_groups[entry["week_num"]].append({
                "day": entry["day_num"], "path": out["path"], "part": out["part"], "of": out["of"],
                "r_start": out["r_start"], "r_end": out["r_end"], "date": entry["date"],
                "pgno": detected_pgno, "keyword": entry["keyword"], "file": entry["file"],
                "page": entry["page"],
            })

        log(f"       ✓ {entry['file']} -> flagged: {entry['highlighted'] or 'none'}")

    for w in week_groups:
        week_groups[w].sort(key=lambda e: (e["day"], e["keyword"], e["page"], e["part"]))

    csv_path = write_csv_report(report_root, csv_results)

    pptx_path = None
    template_path = template_path or manifest.get("template_path")
    if template_path and os.path.exists(template_path) and manifest["entries"]:
        month_name, year = manifest["month_name"], manifest["year"]
        pptx_path = os.path.join(report_root, f"Explorer_SEO_Report_{month_name}_{year}.pptx")
        summary_text = (
            f"This report covers {month_name} {year}. {total_images} Explorer search-result "
            f"screenshot(s) were reviewed across "
            f"{sum(1 for w in week_groups if week_groups[w])} week(s), with {total_highlighted} "
            f"result row(s) flagged for court/legal/financial-crime related mentions requiring "
            f"follow-up."
        )
        build_pptx(template_path, pptx_path, month_name, year, week_groups, summary_text)

    save_manifest(report_root, manifest["entries"], manifest["month_name"], manifest["year"],
                  manifest["month_num"], template_path)

    return total_images, total_highlighted, csv_path, pptx_path


# ════════════════════════════════════════════════════════════════════════
# Main pipeline
# ════════════════════════════════════════════════════════════════════════

def scan_and_detect(drive_root, year, month_name, month_num, template_path=None):
    """
    Phase 1 of the review-before-build flow. Files sit DIRECTLY in each
    day folder — e.g. "Sudesh Rocharam pg 3 no 6,8&10.jpg" — with no
    subfolder involved. A filename containing "no results" (in any
    spacing/case) is skipped entirely, since there's nothing to detect
    or build from it.

    template_path, when given directly, skips the interactive "choose a
    template" prompt entirely — required for any non-Tkinter caller
    (Colab notebooks included), since prompt_choice() blocks waiting on
    a GUI response queue that nothing will ever service outside the
    desktop app's own polling loop.
    """
    log(f"Using folder root: {drive_root}")
    log("=" * 70)
    log("MONTHLY EXPLORER REPORT GENERATOR — detection pass")
    log("=" * 70)
    log(f"\n-> Using Month = {month_name} ({month_num}), Year = {year}")

    candidates = [d for d in list_dirs(drive_root) if MAIN_FOLDER_NAME_HINT.lower() in d.lower() and year in d]
    if len(candidates) == 1:
        year_folder_name = candidates[0]
    elif len(candidates) > 1:
        year_folder_name = prompt_choice("Multiple matching year folders found — which one?", candidates)
    else:
        year_folder_name = prompt_choice("Select the year folder from your Drive root:", list_dirs(drive_root))
    if not year_folder_name:
        raise RuntimeError("No year folder available to use.")
    year_folder_path = os.path.join(drive_root, year_folder_name)

    month_dirs = list_dirs(year_folder_path)
    month_candidates = [d for d in month_dirs if month_name.lower() in d.lower()]
    if len(month_candidates) == 1:
        month_folder_name = month_candidates[0]
    elif len(month_candidates) > 1:
        month_folder_name = prompt_choice("Multiple matching month folders found — which one?", month_candidates)
    else:
        month_folder_name = prompt_choice("Select the month folder:", month_dirs)
    if not month_folder_name:
        raise RuntimeError("No month folder available to use.")
    month_folder_path = os.path.join(year_folder_path, month_folder_name)

    day_dirs_all = list_dirs(month_folder_path)
    log(f"\nMonth folder: '{month_folder_path}'")

    report_root = os.path.join(month_folder_path, "Explorer Report")

    csv_results = defaultdict(lambda: defaultdict(list))
    skipped = []
    jobs = []

    # New, flatter layout: <month>/Explorer/<day>/<file> — files sit
    # directly in each day folder (no further subfolder), with the
    # keyword parsed straight from the filename itself. This replaced
    # the older <month>/<day>/<file> layout specifically so both report
    # types share one consistent shape: <month>/<Google|Explorer>/<day>/.
    log("\nScanning for an 'Explorer' subfolder, then day subfolders inside it ...\n")
    explorer_dir_name = find_subfolder(month_folder_path, "explorer")
    if not explorer_dir_name:
        log("  \u26a0  No 'Explorer' subfolder found directly inside the month folder \u2014 nothing to scan.")
    else:
        explorer_path = os.path.join(month_folder_path, explorer_dir_name)
        day_dirs_parsed = []
        for d in list_dirs(explorer_path):
            day_num = extract_day_number(d)
            if day_num is not None:
                day_dirs_parsed.append((d, day_num, week_for_day(day_num)))
        day_dirs_parsed.sort(key=lambda x: x[1])
        log(f"Found {len(day_dirs_parsed)} day folder(s) recognized inside 'Explorer/'.")

        for day_folder_name, day_num, week_num in day_dirs_parsed:
            day_folder_path = os.path.join(explorer_path, day_folder_name)
            # Files sit directly in the day folder — e.g. "Sudesh Rocharam
            # pg 3 no 6,8&10.jpg" or "Ashley Rocharam no results.jpg"
            # (skipped entirely below) — no further subfolder involved.
            files = list_files(day_folder_path)
            if not files:
                continue

            day_file_count = 0
            for fname in files:
                fname_no_ext = os.path.splitext(fname)[0]
                ext = os.path.splitext(fname)[1].lower()
                filepath = os.path.join(day_folder_path, fname)

                if is_no_results(fname_no_ext):
                    continue
                if is_marked(fname_no_ext):
                    continue

                keyword = match_keyword(fname_no_ext)
                if keyword is None or week_num is None:
                    skipped.append((filepath, "no keyword match or unassignable week"))
                    continue

                if ext not in IMAGE_EXTS:
                    fallback_pgno = None
                    m = re.search(r'PG\s*#?\s*(\d+)\s*(?:NO|#)\s*([\d]+(?:\s*-\s*\d+)?)', fname_no_ext, re.IGNORECASE)
                    if m:
                        fallback_pgno = f"Pg{m.group(1).strip()}#{m.group(2).replace(' ', '')}"
                    if fallback_pgno:
                        csv_results[keyword][week_num].append(fallback_pgno)
                    continue

                date_str = f"{year}-{month_num:02d}-{day_num:02d}"
                jobs.append({
                    "filepath": filepath, "week_num": week_num, "day_num": day_num,
                    "date": date_str, "keyword": keyword,
                })
                day_file_count += 1

            if day_file_count:
                log(f"  Day {day_folder_name} (Week {week_num}): {day_file_count} file(s)")

    log(f"\nScan complete. Detecting rows in {len(jobs)} image(s)...\n")
    set_progress(0, max(len(jobs), 1))

    manifest_entries = []
    done = 0

    def _run_job(job):
        info = detect_rows_for_review(job["filepath"])
        return job, info

    with ThreadPoolExecutor(max_workers=IMAGE_WORKERS) as pool:
        futures = [pool.submit(_run_job, job) for job in jobs]
        for fut in as_completed(futures):
            try:
                job, info = fut.result()
            except Exception as exc:
                log(f"  ✗  detection failed: {exc}")
                continue
            done += 1
            manifest_entries.append({
                "raw_filepath": job["filepath"], "keyword": job["keyword"],
                "week_num": job["week_num"], "day_num": job["day_num"], "date": job["date"],
                "file": os.path.basename(job["filepath"]), "name": strip_pg_suffix(job["keyword"]),
                "page": info["page"], "r_first": info["r_first"], "r_last": info["r_last"],
                "total_rows": info["total_rows"], "x_right": info["x_right"],
                "row_coords": info["row_coords"], "highlighted": info["suggested_flagged"],
                "outputs": [],
            })
            log(f"       ✓ {os.path.basename(job['filepath'])} -> page {info['page']}, "
                f"{info['total_rows']} row(s), suggested flagged: {info['suggested_flagged'] or 'none'}")
            set_progress(done, max(len(jobs), 1))

    log("\nDetection complete.")

    log("\n" + "=" * 70)
    log("POWERPOINT TEMPLATE")
    log("=" * 70)
    if template_path:
        log(f"Using template: {template_path}")
    else:
        pptx_candidates = [f for f in list_files(drive_root) if f.lower().endswith(('.pptx', '.potx'))]
        if not pptx_candidates:
            log("No .pptx template found in the Drive root — CSV/images will still be produced.")
        else:
            template_choice = prompt_choice("Choose a template presentation to match the format/colors of:", pptx_candidates)
            if template_choice:
                template_path = os.path.join(drive_root, template_choice)

    manifest_path = save_manifest(report_root, manifest_entries, month_name, year, month_num, template_path)
    log(f"\nManifest saved to: {manifest_path}")
    if skipped:
        log(f"\n{len(skipped)} file(s) skipped from detection:")
        for filepath, reason in skipped:
            log(f"  - {filepath}\n      reason: {reason}")

    manifest = load_manifest(report_root)
    return report_root, manifest, csv_results


def main(drive_root, year, month_name, month_num):
    report_root, manifest, extra_csv_results = scan_and_detect(drive_root, year, month_name, month_num)

    log("\n" + "=" * 70)
    log("BUILDING REPORT (using auto-detected suggestions)")
    log("=" * 70)
    total_images, total_highlighted, report_path, out_pptx = build_report_from_manifest(
        report_root, manifest, manifest.get("template_path"), extra_csv_results=extra_csv_results)

    log("\n" + "=" * 70)
    log("SUMMARY")
    log("=" * 70)
    log(f"Images processed: {total_images}")
    log(f"Rows flagged: {total_highlighted}")
    log(f"✅ CSV written to: {report_path}")
    if out_pptx:
        log(f"✅ PPTX written to: {out_pptx}")
    log(f"\nTip: run review_report_explorer.py and open '{report_root}' any time to review "
        f"flagged rows visually and rebuild the CSV/PPTX/images with corrections.")
    log("\nDone.")

    return report_path, out_pptx


# ════════════════════════════════════════════════════════════════════════
# GUI
# ════════════════════════════════════════════════════════════════════════

APP_BG = "#1f2430"
PANEL_BG = "#2a3040"
ACCENT = "#4f9dff"
TEXT_FG = "#e8ebf2"
MUTED_FG = "#9aa4b8"


class ChoiceDialog(tk.Toplevel):
    def __init__(self, parent, question, options):
        super().__init__(parent)
        self.title("Please choose")
        self.configure(bg=PANEL_BG)
        self.resizable(False, False)
        self.result = options[0]
        self.protocol("WM_DELETE_WINDOW", lambda: None)
        self.transient(parent)
        self.grab_set()

        tk.Label(self, text=question, bg=PANEL_BG, fg=TEXT_FG, wraplength=420,
                  justify="left", font=("Segoe UI", 10, "bold")).pack(padx=18, pady=(18, 8), anchor="w")

        frame = tk.Frame(self, bg=PANEL_BG)
        frame.pack(padx=18, pady=4, fill="both", expand=True)

        scrollbar = tk.Scrollbar(frame)
        scrollbar.pack(side="right", fill="y")
        self.listbox = tk.Listbox(frame, height=min(8, max(3, len(options))), width=55,
                                   yscrollcommand=scrollbar.set, activestyle="dotbox",
                                   bg="#161a22", fg=TEXT_FG, selectbackground=ACCENT,
                                   font=("Segoe UI", 10))
        for opt in options:
            self.listbox.insert("end", opt)
        self.listbox.selection_set(0)
        self.listbox.pack(side="left", fill="both", expand=True)
        scrollbar.config(command=self.listbox.yview)
        self.listbox.bind("<Double-Button-1>", lambda e: self._confirm())

        btn_frame = tk.Frame(self, bg=PANEL_BG)
        btn_frame.pack(pady=14)
        tk.Button(btn_frame, text="OK", width=12, command=self._confirm,
                  bg=ACCENT, fg="white", relief="flat",
                  activebackground="#3d82e0", font=("Segoe UI", 10, "bold")).pack()

        self.update_idletasks()
        x = parent.winfo_rootx() + (parent.winfo_width() - self.winfo_width()) // 2
        y = parent.winfo_rooty() + (parent.winfo_height() - self.winfo_height()) // 2
        self.geometry(f"+{max(x,0)}+{max(y,0)}")

    def _confirm(self):
        sel = self.listbox.curselection()
        if sel:
            self.result = self.listbox.get(sel[0])
        self.grab_release()
        self.destroy()


class App:
    def __init__(self, root):
        self.root = root
        self.worker_thread = None
        root.title("Explorer Monthly SEO Report Generator")
        root.configure(bg=APP_BG)
        root.geometry("760x640")
        root.minsize(680, 560)

        header = tk.Frame(root, bg=APP_BG)
        header.pack(fill="x", padx=24, pady=(20, 6))
        tk.Label(header, text="Explorer Monthly SEO Report Generator", bg=APP_BG, fg=TEXT_FG,
                  font=("Segoe UI", 18, "bold")).pack(anchor="w")
        tk.Label(header, text="Pick your folder, year, and month — then click Run Report.",
                  bg=APP_BG, fg=MUTED_FG, font=("Segoe UI", 10)).pack(anchor="w", pady=(2, 0))

        panel = tk.Frame(root, bg=PANEL_BG, padx=18, pady=16)
        panel.pack(fill="x", padx=24, pady=10)

        tk.Label(panel, text="1. Folder containing your 'Sudesh <year>' folders",
                  bg=PANEL_BG, fg=TEXT_FG, font=("Segoe UI", 10, "bold")).grid(row=0, column=0, columnspan=3, sticky="w")
        self.folder_var = tk.StringVar(value=os.getcwd())
        folder_entry = tk.Entry(panel, textvariable=self.folder_var, width=52,
                                 bg="#161a22", fg=TEXT_FG, insertbackground=TEXT_FG, relief="flat")
        folder_entry.grid(row=1, column=0, columnspan=2, sticky="we", pady=(4, 14), ipady=4)
        tk.Button(panel, text="Browse…", command=self.browse_folder, bg="#3a4256", fg=TEXT_FG,
                  relief="flat", activebackground="#4a5268").grid(row=1, column=2, padx=(10, 0), pady=(4, 14), sticky="we")

        tk.Label(panel, text="2. Year", bg=PANEL_BG, fg=TEXT_FG,
                  font=("Segoe UI", 10, "bold")).grid(row=2, column=0, sticky="w")
        tk.Label(panel, text="3. Month", bg=PANEL_BG, fg=TEXT_FG,
                  font=("Segoe UI", 10, "bold")).grid(row=2, column=1, sticky="w")

        import datetime
        now = datetime.datetime.now()
        self.year_var = tk.StringVar(value=str(now.year))
        year_spin = tk.Spinbox(panel, from_=2000, to=2100, textvariable=self.year_var, width=10,
                                bg="#161a22", fg=TEXT_FG, insertbackground=TEXT_FG, relief="flat",
                                buttonbackground="#3a4256")
        year_spin.grid(row=3, column=0, sticky="w", pady=(4, 0), ipady=3)

        self.month_var = tk.StringVar(value=MONTH_NAMES[now.month - 1])
        month_combo = ttk.Combobox(panel, textvariable=self.month_var, values=MONTH_NAMES,
                                    state="readonly", width=16)
        month_combo.grid(row=3, column=1, sticky="w", pady=(4, 0))

        panel.grid_columnconfigure(0, weight=1)
        panel.grid_columnconfigure(1, weight=1)
        panel.grid_columnconfigure(2, weight=0)

        run_frame = tk.Frame(root, bg=APP_BG)
        run_frame.pack(fill="x", padx=24, pady=(14, 6))

        self.run_btn = tk.Button(run_frame, text="▶  Run Report", command=self.start_run,
                                  bg=ACCENT, fg="white", font=("Segoe UI", 12, "bold"),
                                  relief="flat", activebackground="#3d82e0", padx=18, pady=8)
        self.run_btn.pack(side="left")

        self.status_var = tk.StringVar(value="Idle — ready when you are.")
        tk.Label(run_frame, textvariable=self.status_var, bg=APP_BG, fg=MUTED_FG,
                  font=("Segoe UI", 10)).pack(side="left", padx=16)

        style = ttk.Style()
        style.theme_use("clam")
        style.configure("Custom.Horizontal.TProgressbar", troughcolor="#161a22",
                         background=ACCENT, thickness=16)
        self.progress = ttk.Progressbar(root, orient="horizontal", mode="determinate",
                                         style="Custom.Horizontal.TProgressbar")
        self.progress.pack(fill="x", padx=24, pady=(0, 12))

        tk.Label(root, text="Activity log", bg=APP_BG, fg=MUTED_FG,
                  font=("Segoe UI", 9, "bold")).pack(anchor="w", padx=24)
        self.log_box = scrolledtext.ScrolledText(root, height=14, bg="#12151c", fg="#c7d0e0",
                                                   insertbackground="#c7d0e0", relief="flat",
                                                   font=("Consolas", 9))
        self.log_box.pack(fill="both", expand=True, padx=24, pady=(4, 20))
        self.log_box.configure(state="disabled")

        self.root.after(120, self._poll)

    def browse_folder(self):
        chosen = filedialog.askdirectory(initialdir=self.folder_var.get() or os.getcwd())
        if chosen:
            self.folder_var.set(chosen)

    def append_log(self, text):
        self.log_box.configure(state="normal")
        self.log_box.insert("end", text + "\n")
        self.log_box.see("end")
        self.log_box.configure(state="disabled")

    def start_run(self):
        drive_root = self.folder_var.get().strip()
        year = self.year_var.get().strip()
        month_name = self.month_var.get().strip()

        if not drive_root or not os.path.isdir(drive_root):
            messagebox.showerror("Missing folder", "Please choose a valid folder first.")
            return
        if not (year.isdigit() and len(year) == 4):
            messagebox.showerror("Invalid year", "Please enter a valid 4-digit year.")
            return
        if month_name not in MONTH_NAMES:
            messagebox.showerror("Invalid month", "Please choose a month from the list.")
            return
        month_num = MONTH_NAMES.index(month_name) + 1

        self.log_box.configure(state="normal")
        self.log_box.delete("1.0", "end")
        self.log_box.configure(state="disabled")
        self.progress["value"] = 0
        self.progress["maximum"] = 100
        self.status_var.set("Working…")
        self.run_btn.config(state="disabled", text="Running…")

        def worker():
            try:
                report_root, manifest, extra_csv_results = scan_and_detect(drive_root, year, month_name, month_num)
                RESULT_QUEUE.put(("ok", report_root))
            except Exception as e:
                log("❌ ERROR: " + str(e))
                log(traceback.format_exc())
                RESULT_QUEUE.put(("error", str(e), None))

        self.worker_thread = threading.Thread(target=worker, daemon=True)
        self.worker_thread.start()

    def _poll(self):
        try:
            while True:
                line = LOG_QUEUE.get_nowait()
                self.append_log(line)
        except queue.Empty:
            pass

        try:
            while True:
                done, total = PROGRESS_QUEUE.get_nowait()
                total = max(total, 1)
                self.progress["maximum"] = total
                self.progress["value"] = done
                self.status_var.set(f"Detecting rows… {done}/{total}")
        except queue.Empty:
            pass

        try:
            while True:
                question, options = CHOICE_REQUEST_QUEUE.get_nowait()
                dlg = ChoiceDialog(self.root, question, options)
                self.root.wait_window(dlg)
                CHOICE_RESPONSE_QUEUE.put(dlg.result)
        except queue.Empty:
            pass

        try:
            while True:
                outcome = RESULT_QUEUE.get_nowait()
                self.run_btn.config(state="normal", text="▶  Run Report")
                if outcome[0] == "ok":
                    _, report_root = outcome
                    self.status_var.set("Detection done ✅ — opening review…")
                    self.progress["value"] = self.progress["maximum"]
                    self._launch_reviewer(report_root)
                else:
                    self.status_var.set("Failed ❌")
                    messagebox.showerror("Something went wrong", outcome[1])
        except queue.Empty:
            pass

        self.root.after(120, self._poll)

    def _launch_reviewer(self, report_root):
        try:
            import review_report_explorer
        except ImportError:
            messagebox.showinfo(
                "Detection complete",
                f"Row detection finished, but review_report_explorer.py wasn't found alongside "
                f"this script, so the review step can't run. Manifest saved to:\n{report_root}",
            )
            return

        review_win = tk.Toplevel(self.root)
        review_report_explorer.ReportReviewApp(review_win, report_root=report_root)


def run_app():
    root = tk.Tk()
    App(root)
    root.mainloop()


if __name__ == "__main__":
    run_app()
